# Copyright 2018-2019 QuantumBlack Visual Analytics Limited
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
# http://www.apache.org/licenses/LICENSE-2.0
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND,
# EXPRESS OR IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES
# OF MERCHANTABILITY, FITNESS FOR A PARTICULAR PURPOSE, AND
# NONINFRINGEMENT. IN NO EVENT WILL THE LICENSOR OR OTHER CONTRIBUTORS
# BE LIABLE FOR ANY CLAIM, DAMAGES, OR OTHER LIABILITY, WHETHER IN AN
# ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM, OUT OF, OR IN
# CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE SOFTWARE.
#
# The QuantumBlack Visual Analytics Limited (“QuantumBlack”) name and logo
# (either separately or in combination, “QuantumBlack Trademarks”) are
# trademarks of QuantumBlack. The License does not grant you any right or
# license to the QuantumBlack Trademarks. You may not use the QuantumBlack
# Trademarks or any confusingly similar mark as a trademark for your product,
#     or use the QuantumBlack Trademarks in any other manner that might cause
# confusion in the marketplace, including but not limited to in advertising,
# on websites, or on software.
#
# See the License for the specific language governing permissions and
# limitations under the License.
"""``PowerPointLocalDataSet`` loads and saves data to McK PPT tables."""
import io
from itertools import cycle
from os.path import dirname, join
from typing import Any, Dict, Optional

import pandas as pd
import pptx
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.slide import Slide
from pptx.table import Table, _Cell, _Row

from kedro.io import AbstractDataSet, DataSetError

__all__ = ["PowerPointLocalDataSet"]


def _format_cells(row: _Row, line: str, like: _Cell):
    for cell, value in zip(row.cells, line.split(",")):
        if like.fill.type == MSO_FILL.BACKGROUND:
            cell.fill.background()
        elif like.fill.type == MSO_FILL.SOLID:
            cell.fill.solid()
            cell.fill.fore_color.rgb = like.fill.fore_color.rgb

        cell.margin_bottom = like.margin_bottom
        cell.margin_left = like.margin_left
        cell.margin_right = like.margin_right
        cell.margin_top = like.margin_top

        run = cell.text_frame.paragraphs[0].add_run()
        run.text = value

        font = like.text_frame.paragraphs[0].runs[0].font
        run.font.color.theme_color = font.color.theme_color
        cell.text_frame._set_font(font.name, font.size.pt, font.bold, font.italic)


class PowerPointLocalDataSet(AbstractDataSet):
    """Load and save data to tables in McK Digital format presentations.

    It uses ``pandas.read_csv`` and ``pandas.DataFrame.to_csv``
    internally, so it supports all allowed ``pandas`` options.

    Example:
    ::

        >>> from kedro.contrib.io.powerpoint import PowerPointLocalDataSet
        >>> import numpy as np
        >>> import pandas as pd
        >>>
        >>> data = pd.DataFrame(np.random.randint(10, size=(10, 4)), columns=["a", "b", "c", "d"])
        >>> data_set = PowerPointLocalDataSet("test.pptx", save_args={"title": "Hello, World!"})
        >>>
        >>> data_set.save(data)
        >>> reloaded = data_set.load()
        >>>
        >>> assert data.equals(reloaded)

    """

    def _describe(self) -> Dict[str, Any]:
        return dict(
            filepath=self._filepath,
            load_args=self._load_args,
            save_args=self._save_args,
        )

    def __init__(
        self,
        filepath: str,
        load_args: Optional[Dict[str, Any]] = None,
        save_args: Optional[Dict[str, Any]] = None,
    ) -> None:
        """Creates a new ``PowerPointLocalDataSet``.

        Args:
            filepath: Either a path to a ``.pptx`` file (a string) or a
                file-like object.  If ``filepath`` is missing or
                ``None``, the built-in default presentation "template"
                is loaded.  For more information, see here:
                https://python-pptx.readthedocs.io/en/latest/api/presentation.html?highlight=presentation#pptx.Presentation
            load_args: Provided to underlying ``pandas.read_csv``
                function.  To find all supported arguments, see here:
                https://pandas.pydata.org/pandas-docs/stable/generated/pandas.read_csv.html
                All defaults are preserved, but ``index_col``, which is
                set to ``0``.

                ``slide_name`` specifies which slide to get.  Available
                cases:
                    * Defaults to ``0``: 1st slide as a ``DataFrame``
                    * ``1``: 2nd slide as a ``DataFrame``
            save_args: Provided to underlying
                ``pandas.DataFrame.to_csv`` function.  To find all
                supported arguments, see here:
                https://pandas.pydata.org/pandas-docs/stable/generated/pandas.DataFrame.to_csv.html
                All defaults are preserved.

                ``title`` replaces "Table" in the title placeholder (if specified).

        Raises:
            DataSetError: When slide ``slide_name`` does not contain a
                table.

        """
        default_load_args = {"slide_name": 0, "index_col": 0}
        default_save_args = {"title": None}
        self._filepath = filepath
        self._load_args = (
            {**default_load_args, **load_args}
            if load_args is not None
            else default_load_args
        )
        self._save_args = (
            {**default_save_args, **save_args}
            if save_args is not None
            else default_save_args
        )

    @staticmethod
    def _get_table(slide: Slide) -> Table:
        for shape in slide.shapes:
            if shape.has_table:
                return shape.table
        raise DataSetError("slide does not contain a table")

    def _load(self) -> pd.DataFrame:
        prs = Presentation(self._filepath)
        load_args = self._load_args.copy()
        slide = prs.slides[load_args.pop("slide_name")]
        table = self._get_table(slide)
        buffer = io.StringIO()
        for row in table.rows:
            buffer.write(",".join(cell.text for cell in row.cells))
            buffer.write("\n")
        buffer.seek(0)
        return pd.read_csv(buffer, **load_args)

    @staticmethod
    def _get_template() -> pptx.presentation.Presentation:
        prs = Presentation(
            join(dirname(__file__), "templates", "Digital_Standard_format_Aug 25.pptx")
        )

        # Delete all slides except for the slide containing the table.
        # See https://github.com/scanny/python-pptx/issues/67#issuecomment-287109303.
        slide_id_list = list(
            enumerate(prs.slides._sldIdLst)  # pylint: disable=protected-access
        )
        del slide_id_list[40]
        for i, slide_id in reversed(slide_id_list):
            prs.part.drop_rel(slide_id.rId)
            del prs.slides._sldIdLst[i]  # pylint: disable=protected-access

        return prs

    def _save(self, data: pd.DataFrame) -> None:
        assert self._save_args.get("index", True)
        prs = self._get_template()
        [slide] = prs.slides
        save_args = self._save_args.copy()
        title = save_args.pop("title")
        if title is not None:
            slide.shapes.title.text = title

        # Replace the example table with an appropriately-sized table.
        # See https://github.com/scanny/python-pptx/issues/246#issuecomment-266124095.
        old_table = self._get_table(slide)._graphic_frame
        nrows, ncols = data.shape
        header_row, first_data_row, second_data_row, *_ = old_table.table.rows
        new_table = slide.shapes.add_table(
            nrows + 1,
            ncols + 1,
            old_table.left,
            old_table.top,
            old_table.width,
            (nrows + 1) * header_row.height,
        )
        old_element = old_table._element
        new_element = new_table._element
        old_element.addnext(new_element)
        old_element.getparent().remove(old_element)

        table = new_table.table
        rows_and_lines = zip(table.rows, data.to_csv(**save_args).splitlines())

        row, line = next(rows_and_lines)
        _format_cells(
            row, line, header_row.cells[1]  # First cell of header row contains no runs
        )

        bands = [first_data_row.cells[0], second_data_row.cells[0]]
        for (row, line), band in zip(rows_and_lines, cycle(bands)):
            _format_cells(row, line, band)

        prs.save(self._filepath)
